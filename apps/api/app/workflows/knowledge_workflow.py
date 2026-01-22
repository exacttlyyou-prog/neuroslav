"""
Workflow индексации документов.
"""
from typing import Dict, Any, List
from loguru import logger
import uuid
from pathlib import Path
import tempfile

from app.services.ollama_service import OllamaService
from app.services.rag_service import RAGService
from app.db.models import KnowledgeItem
from app.db.database import AsyncSessionLocal
from langchain_text_splitters import RecursiveCharacterTextSplitter


class KnowledgeWorkflow:
    """Workflow для индексации документов."""
    
    def __init__(self):
        self.ollama = OllamaService()
        self.rag = RAGService()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
    
    async def process_document(
        self,
        file_content: bytes,
        filename: str,
        file_type: str
    ) -> Dict[str, Any]:
        """
        Обрабатывает документ: извлекает текст, чанкирует, индексирует в Vector DB.
        
        Args:
            file_content: Содержимое файла
            filename: Имя файла
            file_type: MIME тип файла
            
        Returns:
            Словарь с результатом индексации
        """
        try:
            doc_id = f"doc-{uuid.uuid4()}"
            
            # Шаг 1: Извлечение текста в зависимости от типа файла
            text_content = await self._extract_text(file_content, filename, file_type)
            
            if not text_content:
                raise ValueError(f"Не удалось извлечь текст из файла {filename}")
            
            # Шаг 2: Chunking текста
            logger.info(f"Разбиение текста на чанки...")
            chunks = self.text_splitter.split_text(text_content)
            logger.info(f"Создано {len(chunks)} чанков")
            
            # Шаг 3: Индексация каждого чанка в ChromaDB
            chunk_ids = []
            for i, chunk in enumerate(chunks):
                chunk_id = f"{doc_id}-chunk-{i}"
                await self.rag.add_knowledge(
                    doc_id=chunk_id,
                    content=chunk,
                    metadata={
                        "source_file": filename,
                        "file_type": file_type,
                        "chunk_index": i,
                        "total_chunks": len(chunks)
                    }
                )
                chunk_ids.append(chunk_id)
            
            # Шаг 4: Сохранение метаданных в SQLite
            async with AsyncSessionLocal() as session:
                knowledge_item = KnowledgeItem(
                    id=doc_id,
                    source_file=filename,
                    file_type=file_type,
                    chunks_count=len(chunks),
                    metadata={
                        "chunk_ids": chunk_ids,
                        "total_length": len(text_content)
                    }
                )
                session.add(knowledge_item)
                await session.commit()
                await session.refresh(knowledge_item)
            
            logger.info(f"Документ индексирован: {doc_id}")
            
            return {
                "item_id": doc_id,
                "chunks_count": len(chunks),
                "message": f"Документ успешно индексирован ({len(chunks)} чанков)"
            }
            
        except Exception as e:
            logger.error(f"Ошибка при индексации документа: {e}")
            raise
    
    async def _extract_text(self, file_content: bytes, filename: str, file_type: str) -> str:
        """
        Извлекает текст из файла в зависимости от типа.
        
        Args:
            file_content: Содержимое файла
            filename: Имя файла
            file_type: MIME тип файла
            
        Returns:
            Извлеченный текст
        """
        # Сохраняем во временный файл
        with tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix) as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            if file_type == "application/pdf":
                return await self._extract_from_pdf(tmp_path)
            elif file_type in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword"
            ]:
                return await self._extract_from_docx(tmp_path)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return await self._extract_from_pptx(tmp_path)
            else:
                # Пробуем определить по расширению
                ext = Path(filename).suffix.lower()
                if ext == ".pdf":
                    return await self._extract_from_pdf(tmp_path)
                elif ext in [".docx", ".doc"]:
                    return await self._extract_from_docx(tmp_path)
                elif ext == ".pptx":
                    return await self._extract_from_pptx(tmp_path)
                else:
                    raise ValueError(f"Неподдерживаемый тип файла: {file_type}")
        finally:
            # Удаляем временный файл
            Path(tmp_path).unlink(missing_ok=True)
    
    async def _extract_from_pdf(self, file_path: str) -> str:
        """Извлекает текст из PDF."""
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text_parts.append(page.extract_text())
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Ошибка при извлечении текста из PDF: {e}")
            raise
    
    async def _extract_from_docx(self, file_path: str) -> str:
        """Извлекает текст из DOCX."""
        try:
            from docx import Document
            doc = Document(file_path)
            text_parts = []
            for paragraph in doc.paragraphs:
                text_parts.append(paragraph.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Ошибка при извлечении текста из DOCX: {e}")
            raise
    
    async def _extract_from_pptx(self, file_path: str) -> str:
        """Извлекает текст из PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Ошибка при извлечении текста из PPTX: {e}")
            raise
